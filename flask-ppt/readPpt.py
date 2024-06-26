from flask import Flask, jsonify
from pptx import Presentation
from scoreSpeech import scoreText

app = Flask(__name__)

@app.route('/get-scores', methods=['GET'])
def get_scores():
    # Open the presentation
    ppt = Presentation("flask-ppt/testPpt.pptx")
    
    # Extract text from each slide in the presentation
    pptText = []
    for slide in ppt.slides:
        slide_text = ""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    slide_text += run.text
        pptText.append(slide_text)
    
    # Score each extracted text
    scores = []
    for slide_text in pptText:
        scores.append([(scoreText(slide_text, "correctness")), (scoreText(slide_text, "relevance")), (scoreText(slide_text, "clarity")), (scoreText(slide_text, "grammar"))])
    
    return scores

if __name__ == '__main__':
    app.run(debug=True)
